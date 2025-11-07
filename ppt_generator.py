from pptx import Presentation
from pptx.util import Inches

def create_ppt(images, texts, output_filename="generated/presentation.pptx"):
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]

    for i, (image_path, text) in enumerate(zip(images, texts)):
        slide = prs.slides.add_slide(title_only_slide_layout)
        title = slide.shapes.title
        title.text = f"Slide {i+1}"
        
        pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5))
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.text = text

    prs.save(output_filename)
